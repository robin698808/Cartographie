"""
PPTX generator using onepoint template.
Generates slides for the cartographie applicative app.
"""
import io
import math
import datetime
from typing import List, Dict, Any, Optional

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE

try:
    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
    _RECT_TYPE = MSO_AUTO_SHAPE_TYPE.RECTANGLE
except (ImportError, AttributeError):
    # python-pptx >= 1.0 may have moved this
    _RECT_TYPE = 1  # MSO_AUTO_SHAPE_TYPE.RECTANGLE integer value

# ─── CONSTANTS ───────────────────────────────────────────────────────────────
CONTENT_X   = 1.25   # inches - left margin for content
CONTENT_Y   = 2.20   # inches - top of content area (below title)
CONTENT_W   = 10.80  # inches - content width
CONTENT_H   = 4.80   # inches - content height
SLIDE_W     = 13.33  # inches
SLIDE_H     = 7.50   # inches

PV_PAL = [
    "2563EB","475569","0891B2","4F46E5","0E7490","64748B",
    "1D4ED8","0369A1","334155","6366F1","1E40AF","0F766E"
]

D1_COLORS = {
    "Transfert TSA": "F59E0B", "Maintien": "10B981",
    "Rebuild": "6366F1", "Abandon": "EF4444", "Non défini": "94A3B8"
}
D2_COLORS = {
    "Clone & Clean": "3B82F6", "Transfert": "10B981",
    "Rebuild": "8B5CF6", "Abandon": "EF4444", "Non défini": "94A3B8"
}

# ─── COLOR HELPERS ────────────────────────────────────────────────────────────
def hex_color(h: str) -> RGBColor:
    h = h.lstrip("#")
    if len(h) == 3:
        h = "".join(c*2 for c in h)
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

def lighten(h: str, pct: float = 0.7) -> str:
    """Mix color with white by pct (0=original, 1=white)."""
    h = h.lstrip("#")
    r, g, b = int(h[0:2],16), int(h[2:4],16), int(h[4:6],16)
    r2 = int(r + (255-r)*pct)
    g2 = int(g + (255-g)*pct)
    b2 = int(b + (255-b)*pct)
    return f"{r2:02X}{g2:02X}{b2:02X}"

# ─── DRAWING HELPERS ──────────────────────────────────────────────────────────
def add_rect(slide, x, y, w, h, fill_hex=None, line_hex=None, line_w=0.5, transparency=0):
    """Add rectangle. All coords in inches. transparency 0-100."""
    if w <= 0 or h <= 0:
        return None
    shape = slide.shapes.add_shape(
        _RECT_TYPE,
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    if fill_hex:
        shape.fill.solid()
        shape.fill.fore_color.rgb = hex_color(fill_hex)
        if transparency > 0:
            try:
                from lxml import etree
                sp_pr = shape._element.spPr
                ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
                solid_fill = sp_pr.find(f'{{{ns}}}solidFill')
                if solid_fill is not None:
                    srgb = solid_fill.find(f'{{{ns}}}srgbClr')
                    if srgb is not None:
                        alpha = etree.SubElement(srgb, f'{{{ns}}}alpha')
                        alpha.set('val', str(int((100 - transparency) * 1000)))
            except Exception:
                pass  # transparency is cosmetic, skip if it fails
    else:
        shape.fill.background()

    if line_hex:
        shape.line.color.rgb = hex_color(line_hex)
        shape.line.width = Pt(line_w)
    else:
        shape.line.fill.background()

    return shape

def add_text(slide, text, x, y, w, h, font_size=11, bold=False, color_hex="1E293B",
             align="left", italic=False, font_face="Poppins", wrap=True):
    """Add text box. All coords in inches."""
    if w <= 0 or h <= 0:
        return None
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    align_map = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}
    p.alignment = align_map.get(align, PP_ALIGN.LEFT)
    # Remove paragraph spacing
    try:
        from pptx.oxml.ns import qn
        pPr = p._p.get_or_add_pPr()
        pPr.set(qn('a:spcBef'), '0')
        pPr.set(qn('a:spcAft'), '0')
    except Exception:
        pass
    run = p.add_run()
    run.text = str(text)
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    if color_hex:
        run.font.color.rgb = hex_color(color_hex)
    if font_face:
        run.font.name = font_face
    return txBox

def set_placeholders(slide, title=None, breadcrumb=None, subtitle=None):
    """Set placeholder text by partial name matching."""
    for ph in slide.placeholders:
        name = ph.name.lower()
        try:
            if title and ('titre' in name or ('title' in name and 'sous' not in name)):
                ph.text = title
            elif breadcrumb and ('ariane' in name or 'texte 3' in name or 'fil' in name):
                ph.text = breadcrumb
            elif subtitle and ('sous-titre' in name or 'variante' in name or 'texte 2' in name):
                ph.text = subtitle
        except Exception:
            pass

# ─── SLIDE MANAGEMENT ────────────────────────────────────────────────────────
def _clear_slides(prs: Presentation):
    """Remove all example slides while preserving layouts and masters."""
    r_ns = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
    sldIdLst = prs.slides._sldIdLst
    for sldId in list(sldIdLst):
        rId = sldId.get(f'{{{r_ns}}}id')
        if rId:
            try:
                prs.part.drop_rel(rId)
            except Exception:
                pass
        try:
            sldIdLst.remove(sldId)
        except Exception:
            pass

# ─── TREEMAP ALGORITHM ───────────────────────────────────────────────────────
def pv_squarify(values: List[float], rect: Dict) -> List[Dict]:
    """
    Squarified treemap layout. Returns list of {x, y, w, h} rects, one per value.
    Port of the JS pvSquarify function.
    """
    total = sum(v for v in values if v > 0)
    if not total or not values or rect['w'] <= 0 or rect['h'] <= 0:
        return [{'x': rect['x'], 'y': rect['y'], 'w': rect['w'], 'h': rect['h']} for _ in values]

    items = sorted(
        [{'v': v, 'i': i} for i, v in enumerate(values) if v > 0],
        key=lambda x: -x['v']
    )
    results = [{'x': rect['x'], 'y': rect['y'], 'w': rect['w'], 'h': rect['h']} for _ in values]
    cursor = dict(rect)
    remaining = list(items)

    while remaining:
        row = []
        best_ratio = float('inf')
        i = 0
        while i < len(remaining):
            row.append(remaining[i])
            row_sum = sum(r['v'] for r in row)
            row_frac = row_sum / total

            if cursor['w'] >= cursor['h']:
                row_thickness = cursor['h'] * row_frac
                if row_thickness <= 0:
                    break
                max_w = (row[-1]['v'] / row_sum) * cursor['w'] if row_sum > 0 else 0
                ratio = max(row_thickness / max_w, max_w / row_thickness) if max_w > 0 else float('inf')
            else:
                row_thickness = cursor['w'] * row_frac
                if row_thickness <= 0:
                    break
                max_h = (row[-1]['v'] / row_sum) * cursor['h'] if row_sum > 0 else 0
                ratio = max(row_thickness / max_h, max_h / row_thickness) if max_h > 0 else float('inf')

            if ratio > best_ratio:
                row.pop()
                break
            best_ratio = ratio
            i += 1

        if not row:
            row.append(remaining[0])

        row_sum = sum(r['v'] for r in row)
        row_frac = row_sum / total

        if cursor['w'] >= cursor['h']:
            row_thickness = cursor['h'] * row_frac
            cx = cursor['x']
            for item in row:
                ww = (item['v'] / row_sum) * cursor['w']
                results[item['i']] = {'x': cx, 'y': cursor['y'], 'w': ww, 'h': row_thickness}
                cx += ww
            cursor = {
                'x': cursor['x'], 'y': cursor['y'] + row_thickness,
                'w': cursor['w'], 'h': cursor['h'] - row_thickness
            }
        else:
            row_thickness = cursor['w'] * row_frac
            cy = cursor['y']
            for item in row:
                hh = (item['v'] / row_sum) * cursor['h']
                results[item['i']] = {'x': cursor['x'], 'y': cy, 'w': row_thickness, 'h': hh}
                cy += hh
            cursor = {
                'x': cursor['x'] + row_thickness, 'y': cursor['y'],
                'w': cursor['w'] - row_thickness, 'h': cursor['h']
            }

        for item in row:
            remaining.remove(item)

    return results


def pv_build_layout(apps: List[Dict], w: float, h: float) -> List[Dict]:
    """
    3-level treemap: domain → category → app.
    Returns list of domain dicts with nested category/app rects (all in inches).
    w, h: total area dimensions in inches.
    """
    D_HEADER_H = 0.22  # domain header height in inches
    Q_HEADER_H = 0.16  # category header height in inches

    # Group by domain
    dom_map = {}
    for a in apps:
        d = a.get('domain') or 'Autre'
        dom_map.setdefault(d, []).append(a)

    doms_info = sorted(
        [{'domaine': d, 'apps': apps_list, 'nb_apps': len(apps_list)}
         for d, apps_list in dom_map.items()],
        key=lambda x: -x['nb_apps']
    )

    dom_rects = pv_squarify(
        [d['nb_apps'] for d in doms_info],
        {'x': 0, 'y': 0, 'w': w, 'h': h}
    )

    result = []
    for i, d in enumerate(doms_info):
        r = dom_rects[i]
        # Group by category
        cat_map = {}
        for a in d['apps']:
            c = a.get('category') or '—'
            cat_map.setdefault(c, []).append(a)

        cats = sorted(
            [{'quartier': c, 'apps': app_list, 'nb_apps': len(app_list)}
             for c, app_list in cat_map.items()],
            key=lambda x: -x['nb_apps']
        )

        inner = {
            'x': r['x'] + 0.04,
            'y': r['y'] + D_HEADER_H + 0.10,
            'w': max(0, r['w'] - 0.08),
            'h': max(0, r['h'] - D_HEADER_H - 0.14)
        }

        q_rects = pv_squarify([c['nb_apps'] for c in cats], inner)

        quartiers = []
        for j, cat in enumerate(cats):
            qr = q_rects[j]
            app_inner = {
                'x': qr['x'] + 0.02,
                'y': qr['y'] + Q_HEADER_H + 0.02,
                'w': max(0, qr['w'] - 0.04),
                'h': max(0, qr['h'] - Q_HEADER_H - 0.04)
            }
            app_rects = pv_squarify([1.0] * cat['nb_apps'], app_inner)

            quartiers.append({
                'quartier': cat['quartier'],
                'nb_apps': cat['nb_apps'],
                'rect': qr,
                'apps': [{'app': a, 'rect': app_rects[k]} for k, a in enumerate(cat['apps'])]
            })

        result.append({
            'domaine': d['domaine'],
            'nb_apps': d['nb_apps'],
            'rect': r,
            'quartiers': quartiers
        })

    return result


# ─── SLIDE GENERATORS ─────────────────────────────────────────────────────────

def add_cover_slide(prs: Presentation, apps: List[Dict], flows: List[Dict], opts: Dict):
    """Slide 1: Cover / Page de garde."""
    layout = prs.slide_layouts[0]  # 1_Page de garde_Light
    slide = prs.slides.add_slide(layout)

    # Try to set placeholders
    today = datetime.date.today().strftime("%d/%m/%Y")
    client_name = opts.get('clientName', '')

    for ph in slide.placeholders:
        name = ph.name.lower()
        try:
            if 'titre' in name and 'sous' not in name:
                ph.text = "Cartographie Applicative"
            elif 'sous-titre' in name or 'description' in name:
                ph.text = client_name or "Analyse du Système d'Information"
            elif 'date' in name:
                ph.text = today
        except Exception:
            pass


def add_synthese_messages_slide(prs: Presentation, apps: List[Dict], flows: List[Dict], opts: Dict):
    """Slide 2: Synthèse & Messages clés."""
    layout = prs.slide_layouts[8]  # 4_Contenu par défaut
    slide = prs.slides.add_slide(layout)
    set_placeholders(slide,
        title="SYNTHÈSE — MESSAGES CLÉS",
        breadcrumb="Cartographie Applicative",
        subtitle="Concentration des flux · Domaines · Hubs"
    )

    primary = opts.get('clientPrimary', '2979FF')

    # Compute data
    # Pairs
    pairs = {}
    for f in flows:
        fa = next((a for a in apps if a['id'] == f.get('from')), None)
        ta = next((a for a in apps if a['id'] == f.get('to')), None)
        if fa and ta and fa.get('domain') != ta.get('domain'):
            k = f"{fa['domain']} → {ta['domain']}"
            pairs[k] = pairs.get(k, 0) + 1
    top_pairs = sorted(pairs.items(), key=lambda x: -x[1])[:5]
    tot_inter = sum(pairs.values())

    # Domain flows
    dom_flows = {}
    for a in apps:
        dom_flows[a.get('domain', 'Autre')] = 0
    for f in flows:
        fa = next((a for a in apps if a['id'] == f.get('from')), None)
        ta = next((a for a in apps if a['id'] == f.get('to')), None)
        if fa:
            dom_flows[fa.get('domain', 'Autre')] = dom_flows.get(fa.get('domain', 'Autre'), 0) + 1
        if ta:
            dom_flows[ta.get('domain', 'Autre')] = dom_flows.get(ta.get('domain', 'Autre'), 0) + 1
    top_dom_flows = sorted(dom_flows.items(), key=lambda x: -x[1])[:5]

    # Hubs
    conn = {a['id']: 0 for a in apps}
    for f in flows:
        if f.get('from') in conn:
            conn[f['from']] += 1
        if f.get('to') in conn:
            conn[f['to']] += 1
    top_hubs = sorted([a for a in apps if conn.get(a['id'], 0) > 0],
                      key=lambda a: -conn.get(a['id'], 0))[:5]

    # 3 cards
    card_w = 3.20
    card_h = 4.00
    card_gap = 0.20
    card_x0 = CONTENT_X
    card_y0 = CONTENT_Y

    cards = [
        {'title': 'CONCENTRATION DES FLUX', 'accent': '6366F1', 'type': 'pairs'},
        {'title': 'DOMAINES PAR VOLUME DE FLUX', 'accent': '0EA5E9', 'type': 'domflow'},
        {'title': "HUBS DU SYSTÈME D'INFORMATION", 'accent': 'F59E0B', 'type': 'hubs'},
    ]

    for ci, card in enumerate(cards):
        cx = card_x0 + ci * (card_w + card_gap)
        cy = card_y0
        acc = card['accent']

        # Card background
        add_rect(slide, cx, cy, card_w, card_h, fill_hex='FFFFFF', line_hex='E2E8F0', line_w=0.5)
        # Accent top bar
        add_rect(slide, cx, cy, card_w, 0.28, fill_hex=acc, transparency=85)
        add_rect(slide, cx, cy, 0.06, card_h, fill_hex=acc)
        # Title
        add_text(slide, card['title'], cx+0.12, cy+0.06, card_w-0.20, 0.20,
                 font_size=7, bold=True, color_hex=acc)

        # Divider
        add_rect(slide, cx+0.12, cy+0.30, card_w-0.24, 0.01, fill_hex='E2E8F0')

        if card['type'] == 'pairs':
            if not top_pairs:
                add_text(slide, "Aucun flux inter-domaine", cx+0.12, cy+0.50, card_w-0.24, 0.30,
                         font_size=9, color_hex='94A3B8')
            else:
                max_v = top_pairs[0][1]
                for pi, (pair_k, pair_v) in enumerate(top_pairs):
                    py = cy + 0.44 + pi * 0.68
                    if py + 0.60 > cy + card_h:
                        break
                    pct = round(pair_v / tot_inter * 100) if tot_inter else 0
                    add_text(slide, str(pair_v), cx+0.12, py, 0.45, 0.30,
                             font_size=20, bold=True, color_hex=acc)
                    add_text(slide, f"flux — {pct}%", cx+0.60, py+0.12, card_w-0.80, 0.18,
                             font_size=7, color_hex='64748B')
                    add_text(slide, pair_k, cx+0.12, py+0.32, card_w-0.24, 0.18,
                             font_size=8, color_hex='334155')
                    bw = (pair_v / max_v) * (card_w - 0.30) if max_v else 0
                    add_rect(slide, cx+0.12, py+0.50, bw, 0.05, fill_hex=acc)

        elif card['type'] == 'domflow':
            if not top_dom_flows:
                add_text(slide, "Aucun flux défini", cx+0.12, cy+0.50, card_w-0.24, 0.30,
                         font_size=9, color_hex='94A3B8')
            else:
                max_v = top_dom_flows[0][1] if top_dom_flows else 1
                for di, (dom, dv) in enumerate(top_dom_flows):
                    dy = cy + 0.44 + di * 0.68
                    if dy + 0.50 > cy + card_h:
                        break
                    bw = max(0.05, (dv / max_v) * (card_w - 0.70)) if max_v else 0.05
                    add_text(slide, dom, cx+0.12, dy, card_w-0.36, 0.22,
                             font_size=9, bold=True, color_hex='0F172A')
                    add_rect(slide, cx+0.12, dy+0.26, card_w-0.28, 0.12, fill_hex='F1F5F9')
                    add_rect(slide, cx+0.12, dy+0.26, bw, 0.12, fill_hex=acc)
                    add_text(slide, f"{dv} flux", cx+0.18+bw, dy+0.24, card_w-0.36-bw, 0.16,
                             font_size=7.5, bold=True, color_hex=acc)

        elif card['type'] == 'hubs':
            if not top_hubs:
                add_text(slide, "Aucun flux défini", cx+0.12, cy+0.50, card_w-0.24, 0.30,
                         font_size=9, color_hex='94A3B8')
            else:
                max_conn = conn.get(top_hubs[0]['id'], 1) or 1
                for hi, hub in enumerate(top_hubs):
                    hy = cy + 0.44 + hi * 0.68
                    if hy + 0.50 > cy + card_h:
                        break
                    hub_conn = conn.get(hub['id'], 0)
                    bw = (hub_conn / max_conn) * (card_w - 0.70) if max_conn else 0
                    add_text(slide, hub.get('name', ''), cx+0.12, hy, card_w-0.28, 0.22,
                             font_size=9, bold=True, color_hex='0F172A')
                    add_rect(slide, cx+0.12, hy+0.26, bw, 0.12, fill_hex=acc)
                    add_text(slide, f"{hub_conn} cx", cx+0.18+bw, hy+0.24, 0.50, 0.16,
                             font_size=8, bold=True, color_hex=acc)
                    add_text(slide, hub.get('domain', ''), cx+0.12, hy+0.42, card_w-0.24, 0.14,
                             font_size=7, color_hex='64748B')

    # KPI band at bottom
    kpi_y = card_y0 + card_h + 0.12
    kpi_h = 0.60
    add_rect(slide, CONTENT_X, kpi_y, CONTENT_W, kpi_h, fill_hex='FFFFFF', line_hex='E2E8F0', line_w=0.5)

    kpis = [
        {'l': 'Total flux', 'v': len(flows), 'c': '6366F1'},
        {'l': 'Flux inter-domaines', 'v': tot_inter, 'c': '8B5CF6'},
        {'l': 'Ratio flux/app', 'v': round(len(flows)/len(apps), 1) if apps else 0, 'c': '22D3EE'},
    ]
    kpi_col_w = CONTENT_W / len(kpis)
    for ki, kpi in enumerate(kpis):
        kx = CONTENT_X + ki * kpi_col_w
        add_rect(slide, kx + 0.08, kpi_y + 0.06, 0.04, kpi_h - 0.12, fill_hex=kpi['c'])
        add_text(slide, str(kpi['v']), kx + 0.18, kpi_y + 0.06, kpi_col_w - 0.30, 0.34,
                 font_size=22, bold=True, color_hex=kpi['c'])
        add_text(slide, kpi['l'], kx + 0.18, kpi_y + 0.38, kpi_col_w - 0.30, 0.18,
                 font_size=8, color_hex='64748B')


def add_executive_slide(prs: Presentation, apps: List[Dict], opts: Dict):
    """Slide 3: Synthèse exécutive — KPI + D1/D2 donuts + owners."""
    layout = prs.slide_layouts[8]
    slide = prs.slides.add_slide(layout)
    set_placeholders(slide,
        title="SYNTHÈSE EXÉCUTIVE",
        breadcrumb="Cartographie Applicative",
    )

    primary = opts.get('clientPrimary', '2979FF')

    # KPI tiles
    doms = list({a.get('domain', 'Autre') for a in apps})
    kpis = [
        {'l': 'Applications', 'v': str(len(apps)), 'c': '6366F1'},
        {'l': 'Domaines', 'v': str(len(doms)), 'c': '8B5CF6'},
    ]
    kpi_w = 2.0
    kpi_h = 0.90
    kpi_y = CONTENT_Y
    for ki, kpi in enumerate(kpis):
        kx = CONTENT_X + ki * (kpi_w + 0.20)
        add_rect(slide, kx, kpi_y, kpi_w, kpi_h, fill_hex='FFFFFF', line_hex='E2E8F0', line_w=0.5)
        add_rect(slide, kx, kpi_y, 0.055, kpi_h, fill_hex=kpi['c'])
        add_text(slide, kpi['v'], kx+0.14, kpi_y+0.06, kpi_w-0.20, 0.50,
                 font_size=28, bold=True, color_hex=kpi['c'])
        add_text(slide, kpi['l'], kx+0.14, kpi_y+0.62, kpi_w-0.20, 0.22,
                 font_size=8, color_hex='64748B')

    # D1 donut chart
    d1_stats = [
        {'l': 'Transfert TSA', 'v': sum(1 for a in apps if a.get('statusD1') == 'Transfert TSA'), 'c': 'F59E0B'},
        {'l': 'Maintien', 'v': sum(1 for a in apps if a.get('statusD1') == 'Maintien'), 'c': '10B981'},
        {'l': 'Rebuild', 'v': sum(1 for a in apps if a.get('statusD1') == 'Rebuild'), 'c': '6366F1'},
        {'l': 'Abandon', 'v': sum(1 for a in apps if a.get('statusD1') == 'Abandon'), 'c': 'EF4444'},
        {'l': 'Non défini', 'v': sum(1 for a in apps if not a.get('statusD1')), 'c': '94A3B8'},
    ]
    d2_stats = [
        {'l': 'Clone & Clean', 'v': sum(1 for a in apps if a.get('statusD2') == 'Clone & Clean'), 'c': '3B82F6'},
        {'l': 'Transfert', 'v': sum(1 for a in apps if a.get('statusD2') == 'Transfert'), 'c': '10B981'},
        {'l': 'Rebuild', 'v': sum(1 for a in apps if a.get('statusD2') == 'Rebuild'), 'c': '8B5CF6'},
        {'l': 'Abandon', 'v': sum(1 for a in apps if a.get('statusD2') == 'Abandon'), 'c': 'EF4444'},
        {'l': 'Non défini', 'v': sum(1 for a in apps if not a.get('statusD2')), 'c': '94A3B8'},
    ]

    donut_y = kpi_y + kpi_h + 0.20
    donut_h = 2.60
    donut_w = 4.80

    def draw_donut_block(bx, by, bw, bh, title, title_color, stats):
        add_rect(slide, bx, by, bw, bh, fill_hex='FFFFFF', line_hex='E2E8F0', line_w=0.5)
        add_rect(slide, bx, by, bw, 0.26, fill_hex=lighten(title_color, 0.85))
        add_text(slide, title, bx+0.12, by+0.04, bw-0.20, 0.20,
                 font_size=8, bold=True, color_hex=title_color)

        # Chart data — only non-zero values
        chart_stats = [s for s in stats if s['v'] > 0]
        if chart_stats and len(apps) > 0:
            chart_data = ChartData()
            chart_data.categories = [s['l'] for s in chart_stats]
            chart_data.add_series('', [s['v'] for s in chart_stats])

            chart_sz = min(bh - 0.40, bw * 0.40)
            chart = slide.shapes.add_chart(
                XL_CHART_TYPE.DOUGHNUT,
                Inches(bx + 0.10), Inches(by + 0.30),
                Inches(chart_sz), Inches(chart_sz),
                chart_data
            )
            ch = chart.chart
            ch.has_legend = False
            if hasattr(ch, 'series'):
                for si, ser in enumerate(ch.series):
                    for pi, pt in enumerate(ser.points):
                        try:
                            pt.format.fill.solid()
                            pt.format.fill.fore_color.rgb = hex_color(chart_stats[pi]['c'])
                        except Exception:
                            pass

        # Legend
        leg_x = bx + bw * 0.42
        leg_y = by + 0.32
        row_h = min(0.36, (bh - 0.40) / max(len(stats), 1))
        for s in stats:
            pct = round(s['v'] / len(apps) * 100) if apps else 0
            col = s['c'] if s['v'] > 0 else 'D1D5DB'
            add_rect(slide, leg_x, leg_y + (row_h - 0.12) / 2, 0.12, 0.12, fill_hex=col)
            add_text(slide, s['l'], leg_x + 0.17, leg_y, bw - (leg_x - bx) - 0.70, row_h,
                     font_size=7.5, color_hex='374151' if s['v'] > 0 else '9CA3AF')
            add_text(slide, f"{s['v']} | {pct}%", bx + bw - 0.70, leg_y, 0.65, row_h,
                     font_size=7, color_hex='94A3B8', align='right')
            leg_y += row_h

    draw_donut_block(CONTENT_X, donut_y, donut_w, donut_h,
                     "VISION CLOSING — DAY 1", "D97706", d1_stats)
    draw_donut_block(CONTENT_X + donut_w + 0.25, donut_y, donut_w, donut_h,
                     "VISION CIBLE — DAY 2", "7C3AED", d2_stats)

    # Owner bar chart
    owners = {}
    for a in apps:
        o = (a.get('owner') or '').strip() or 'Non renseigné'
        owners[o] = owners.get(o, 0) + 1
    top_owners = sorted([(k, v) for k, v in owners.items() if k != 'Non renseigné'],
                        key=lambda x: -x[1])[:6]

    ow_y = donut_y + donut_h + 0.15
    ow_h = 1.10
    add_rect(slide, CONTENT_X, ow_y, CONTENT_W, ow_h, fill_hex='FFFFFF', line_hex='E2E8F0', line_w=0.5)
    add_text(slide, "RÉPARTITION PAR RESPONSABLE", CONTENT_X + 0.15, ow_y + 0.06,
             5.0, 0.18, font_size=8, bold=True, color_hex=primary)

    if top_owners:
        max_v = top_owners[0][1]
        col_w = CONTENT_W / max(len(top_owners), 1)
        for i, (owner, cnt) in enumerate(top_owners):
            ox = CONTENT_X + i * col_w + 0.10
            bw_bar = (cnt / max_v) * (col_w - 0.20) if max_v else 0
            add_rect(slide, ox, ow_y + 0.30, bw_bar, 0.18, fill_hex=primary)
            add_text(slide, str(cnt), ox + bw_bar + 0.04, ow_y + 0.29, 0.30, 0.20,
                     font_size=8, bold=True, color_hex=primary)
            add_text(slide, owner, ox, ow_y + 0.52, col_w - 0.20, 0.18,
                     font_size=7, color_hex='475569')
            add_text(slide, f"{round(cnt/len(apps)*100)}%", ox, ow_y + 0.72, col_w - 0.20, 0.16,
                     font_size=6.5, color_hex='94A3B8')


def add_domain_status_slides(prs: Presentation, apps: List[Dict], opts: Dict):
    """Slides: Cartographie par domaine — Day 1 & Day 2 (chip grid)."""
    primary = opts.get('clientPrimary', '2979FF')

    configs = [
        {'field': 'statusD1', 'title': 'CARTOGRAPHIE PAR DOMAINE — DAY 1 (CLOSING)',
         'color_map': D1_COLORS},
        {'field': 'statusD2', 'title': 'CARTOGRAPHIE PAR DOMAINE — DAY 2 (CIBLE)',
         'color_map': D2_COLORS},
    ]

    for cfg in configs:
        doms = list(dict.fromkeys(a.get('domain', 'Autre') for a in apps))
        chip_w, chip_h = 0.90, 0.24
        chip_gap_h, chip_gap_v = 0.04, 0.03
        label_w = 1.80
        content_x2 = CONTENT_X + label_w + 0.10
        chips_area_w = CONTENT_W - label_w - 0.10
        chips_per_row = max(1, int((chips_area_w + chip_gap_h) / (chip_w + chip_gap_h)))

        def make_slide():
            layout = prs.slide_layouts[8]
            sl = prs.slides.add_slide(layout)
            set_placeholders(sl, title=cfg['title'], breadcrumb="Cartographie Applicative")
            # Legend
            leg_x = CONTENT_X
            leg_y = CONTENT_Y - 0.02
            for st, col in cfg['color_map'].items():
                add_rect(sl, leg_x, leg_y + 0.04, 0.14, 0.12, fill_hex=col)
                add_text(sl, st, leg_x + 0.18, leg_y, 1.20, 0.22, font_size=7.5, color_hex='444444')
                leg_x += 1.55
            return sl

        cur_slide = make_slide()
        cur_y = CONTENT_Y + 0.28  # below legend
        max_y = SLIDE_H - 0.50
        dom_gap = 0.06
        dom_pad = 0.05

        for dom in doms:
            dom_apps = [a for a in apps if a.get('domain') == dom]
            chip_rows = math.ceil(len(dom_apps) / chips_per_row) or 1
            dom_h = dom_pad * 2 + chip_rows * chip_h + (chip_rows - 1) * chip_gap_v

            if cur_y + dom_h > max_y:
                cur_slide = make_slide()
                cur_y = CONTENT_Y + 0.28

            # Domain label
            add_rect(cur_slide, CONTENT_X, cur_y, label_w, dom_h, fill_hex='EFF6FF', line_hex='BFDBFE', line_w=0.5)
            add_text(cur_slide, dom, CONTENT_X + 0.08, cur_y + dom_pad,
                     label_w - 0.12, dom_h * 0.55, font_size=8, bold=True, color_hex='1D4ED8')
            add_text(cur_slide, f"{len(dom_apps)} app{'s' if len(dom_apps) > 1 else ''}",
                     CONTENT_X + 0.08, cur_y + dom_h * 0.62,
                     label_w - 0.12, dom_h * 0.32, font_size=6, color_hex='888888')

            # Chips
            for ai, app in enumerate(dom_apps):
                row = ai // chips_per_row
                col = ai % chips_per_row
                chip_x = content_x2 + col * (chip_w + chip_gap_h)
                chip_y = cur_y + dom_pad + row * (chip_h + chip_gap_v)
                st = app.get(cfg['field']) or 'Non défini'
                stc = cfg['color_map'].get(st, '94A3B8')
                add_rect(cur_slide, chip_x, chip_y, chip_w, chip_h,
                         fill_hex=stc, transparency=82, line_hex=stc, line_w=0.5)
                add_rect(cur_slide, chip_x, chip_y, 0.04, chip_h, fill_hex=stc)
                add_text(cur_slide, app.get('name', ''), chip_x + 0.07, chip_y,
                         chip_w - 0.10, chip_h * 0.55, font_size=5.5, bold=True, color_hex='1a1a1a')
                add_text(cur_slide, '—' if st == 'Non défini' else st,
                         chip_x + 0.07, chip_y + chip_h * 0.55,
                         chip_w - 0.10, chip_h * 0.42, font_size=5, color_hex=stc)

            add_rect(cur_slide, CONTENT_X, cur_y + dom_h, CONTENT_W, 0.01, fill_hex='E2E8F0')
            cur_y += dom_h + dom_gap


def add_paysage_slide(prs: Presentation, apps: List[Dict]):
    """Paysage applicatif treemap slide."""
    layout = prs.slide_layouts[9]  # 4_Contenu_pour illustration
    slide = prs.slides.add_slide(layout)
    set_placeholders(slide,
        title="PAYSAGE APPLICATIF",
        breadcrumb="Cartographie Applicative",
        subtitle="Surface proportionnelle au nombre d'applications par domaine"
    )

    # Content area for treemap
    tm_x = CONTENT_X
    tm_y = CONTENT_Y
    tm_w = CONTENT_W
    tm_h = CONTENT_H

    pv_layout = pv_build_layout(apps, tm_w, tm_h)

    for di, dom in enumerate(pv_layout):
        color = PV_PAL[di % len(PV_PAL)]
        r = dom['rect']

        if r['w'] < 0.05 or r['h'] < 0.05:
            continue

        # Domain background
        add_rect(slide,
                 tm_x + r['x'], tm_y + r['y'], r['w'], r['h'],
                 fill_hex=color, transparency=95, line_hex=color, line_w=1.0)

        # Domain header
        hdr_h = 0.22
        add_rect(slide,
                 tm_x + r['x'] + 0.02, tm_y + r['y'] + 0.02,
                 max(0, r['w'] - 0.04), hdr_h,
                 fill_hex=color, transparency=22)

        # Domain name
        if r['w'] > 0.60:
            add_text(slide, dom['domaine'],
                     tm_x + r['x'] + 0.06, tm_y + r['y'] + 0.03,
                     max(0.10, r['w'] - 0.50), hdr_h - 0.04,
                     font_size=max(5.0, min(9.0, r['w'] * 4)),
                     bold=True, color_hex=color)
            add_text(slide, str(dom['nb_apps']),
                     tm_x + r['x'] + r['w'] - 0.50, tm_y + r['y'] + 0.03,
                     0.45, hdr_h - 0.04,
                     font_size=max(5.0, min(8.0, r['w'] * 3.5)),
                     bold=False, color_hex='475569', align='right')

        # Categories
        for cat in dom['quartiers']:
            qr = cat['rect']
            if qr['w'] < 0.04 or qr['h'] < 0.04:
                continue

            # Category background
            add_rect(slide,
                     tm_x + qr['x'] + 0.01, tm_y + qr['y'] + 0.01,
                     max(0, qr['w'] - 0.02), max(0, qr['h'] - 0.02),
                     fill_hex='FFFFFF', transparency=0, line_hex=color, line_w=0.5)

            # Category label
            q_hdr_h = 0.16
            if qr['h'] > 0.20 and qr['w'] > 0.40:
                add_text(slide, cat['quartier'],
                         tm_x + qr['x'] + 0.04, tm_y + qr['y'] + 0.02,
                         max(0.10, qr['w'] - 0.08), q_hdr_h - 0.02,
                         font_size=max(4.5, min(8.0, qr['w'] * 3.5)),
                         bold=True, color_hex=color)

            # Apps
            for item in cat['apps']:
                ar = item['rect']
                if ar['w'] < 0.02 or ar['h'] < 0.02:
                    continue

                add_rect(slide,
                         tm_x + ar['x'], tm_y + ar['y'], ar['w'], ar['h'],
                         fill_hex=color, transparency=90, line_hex=color, line_w=0.3)

                # App name (only if big enough)
                if ar['w'] > 0.55 and ar['h'] > 0.18:
                    add_text(slide, item['app'].get('name', ''),
                             tm_x + ar['x'] + 0.03, tm_y + ar['y'] + 0.02,
                             max(0.10, ar['w'] - 0.06), ar['h'] - 0.04,
                             font_size=max(4.0, min(7.0, ar['w'] * 5)),
                             bold=False, color_hex='334155')


def add_recap_apps_slides(prs: Presentation, apps: List[Dict], opts: Dict):
    """Application recap table, paginated."""
    primary = opts.get('clientPrimary', '2979FF')
    rows_per_slide = 22

    # Sort apps by domain then name
    sorted_apps = sorted(apps, key=lambda a: (a.get('domain', ''), a.get('name', '')))

    headers = ['Nom', 'Domaine', 'Catégorie', 'Criticité', 'Day 1', 'Day 2', 'Éditeur', 'Responsable']
    col_widths = [2.20, 1.40, 1.40, 0.90, 1.10, 1.10, 1.20, 1.00]  # sum ~10.30
    total_w = sum(col_widths)

    def make_recap_slide():
        layout = prs.slide_layouts[8]
        sl = prs.slides.add_slide(layout)
        set_placeholders(sl, title="RÉCAPITULATIF DES APPLICATIONS", breadcrumb="Cartographie Applicative")
        return sl

    for page_start in range(0, max(1, len(sorted_apps)), rows_per_slide):
        page_apps = sorted_apps[page_start:page_start + rows_per_slide]
        slide = make_recap_slide()

        row_h = 0.20
        hdr_h = 0.24
        table_y = CONTENT_Y
        table_x = CONTENT_X

        # Header row
        cx = table_x
        for ci, (col, cw) in enumerate(zip(headers, col_widths)):
            add_rect(slide, cx, table_y, cw, hdr_h, fill_hex=primary)
            add_text(slide, col, cx + 0.04, table_y + 0.02, cw - 0.06, hdr_h - 0.04,
                     font_size=7, bold=True, color_hex='FFFFFF')
            cx += cw

        # Data rows
        for ri, app in enumerate(page_apps):
            row_y = table_y + hdr_h + ri * row_h
            bg = 'F8FAFC' if ri % 2 == 0 else 'FFFFFF'
            cx = table_x
            cells = [
                app.get('name', ''),
                app.get('domain', ''),
                app.get('category', ''),
                app.get('criticality', ''),
                app.get('statusD1', ''),
                app.get('statusD2', ''),
                app.get('vendor', ''),
                app.get('owner', ''),
            ]
            for ci, (cell, cw) in enumerate(zip(cells, col_widths)):
                add_rect(slide, cx, row_y, cw, row_h, fill_hex=bg, line_hex='E2E8F0', line_w=0.3)
                color = primary if ci == 0 else '334155'
                add_text(slide, str(cell or ''), cx + 0.04, row_y + 0.01, cw - 0.06, row_h - 0.02,
                         font_size=6.5, bold=(ci == 0), color_hex=color)
                cx += cw


def add_flux_table_slides(prs: Presentation, apps: List[Dict], flows: List[Dict], opts: Dict):
    """Flux table slides, paginated."""
    if not flows:
        return

    primary = opts.get('clientPrimary', '2979FF')
    rows_per_slide = 25

    # Sort flows
    def flow_key(f):
        fa = next((a for a in apps if a['id'] == f.get('from')), {})
        ta = next((a for a in apps if a['id'] == f.get('to')), {})
        return (fa.get('domain', ''), fa.get('name', ''), ta.get('name', ''))

    sorted_flows = sorted(flows, key=flow_key)

    headers = ['N°', 'Émetteur', 'Récepteur', 'Nom du flux', 'Protocole', 'Fréquence']
    col_widths = [0.40, 2.20, 2.20, 3.40, 1.30, 1.30]

    def make_flux_slide():
        layout = prs.slide_layouts[8]
        sl = prs.slides.add_slide(layout)
        set_placeholders(sl, title="CARTOGRAPHIE DES FLUX", breadcrumb="Cartographie Applicative",
                         subtitle=f"Données : {len(apps)} applications · {len(flows)} flux")
        return sl

    app_by_id = {a['id']: a for a in apps}

    for page_start in range(0, max(1, len(sorted_flows)), rows_per_slide):
        page_flows = sorted_flows[page_start:page_start + rows_per_slide]
        slide = make_flux_slide()

        row_h = 0.18
        hdr_h = 0.22
        table_y = CONTENT_Y
        table_x = CONTENT_X

        # Header
        cx = table_x
        for col, cw in zip(headers, col_widths):
            add_rect(slide, cx, table_y, cw, hdr_h, fill_hex=primary)
            add_text(slide, col, cx + 0.04, table_y + 0.02, cw - 0.06, hdr_h - 0.04,
                     font_size=7, bold=True, color_hex='FFFFFF')
            cx += cw

        # Rows
        for ri, flow in enumerate(page_flows):
            row_y = table_y + hdr_h + ri * row_h
            bg = 'F8FAFC' if ri % 2 == 0 else 'FFFFFF'
            fa = app_by_id.get(flow.get('from', ''), {})
            ta = app_by_id.get(flow.get('to', ''), {})
            cells = [
                str(page_start + ri + 1),
                fa.get('name', ''),
                ta.get('name', ''),
                flow.get('label', '') or flow.get('name', ''),
                flow.get('protocol', ''),
                flow.get('frequency', ''),
            ]
            cx = table_x
            for ci, (cell, cw) in enumerate(zip(cells, col_widths)):
                add_rect(slide, cx, row_y, cw, row_h, fill_hex=bg, line_hex='E2E8F0', line_w=0.3)
                add_text(slide, str(cell or ''), cx + 0.03, row_y + 0.01, cw - 0.05, row_h - 0.02,
                         font_size=6.5, color_hex='334155' if ci > 0 else '1E293B', bold=(ci == 0))
                cx += cw


def add_flux_matrix_slide(prs: Presentation, apps: List[Dict], flows: List[Dict], opts: Dict):
    """Slide : matrice flux inter-domaines (domaine × domaine)."""
    if not flows:
        return
    primary = opts.get('clientPrimary', '2979FF')
    layout = prs.slide_layouts[8]
    slide = prs.slides.add_slide(layout)
    set_placeholders(slide, title="FLUX — VUE AGRÉGÉE PAR DOMAINE",
                     breadcrumb="Cartographie Applicative",
                     subtitle=f"{len(flows)} flux · {len(apps)} applications")

    app_by_id = {a['id']: a for a in apps}
    doms = list(dict.fromkeys(a.get('domain', 'Autre') for a in apps))

    # Build matrix
    matrix = {d: {d2: 0 for d2 in doms} for d in doms}
    for f in flows:
        fa = app_by_id.get(f.get('from', ''), {})
        ta = app_by_id.get(f.get('to', ''), {})
        fd, td = fa.get('domain'), ta.get('domain')
        if fd and td and fd != td:
            matrix[fd][td] = matrix[fd].get(td, 0) + 1

    # Max value for colour scaling
    max_v = max((matrix[r][c] for r in doms for c in doms), default=1) or 1

    # Sizes
    n = len(doms)
    cell_w = min(0.95, (CONTENT_W - 1.60) / n)
    cell_h = min(0.38, (CONTENT_H - 0.40) / n)
    lbl_w = max(1.20, CONTENT_W - n * cell_w - 0.10)
    lbl_h = cell_h
    hdr_h = lbl_h
    tbl_x = CONTENT_X
    tbl_y = CONTENT_Y

    # Column headers (domain names)
    for ci, dom in enumerate(doms):
        cx = tbl_x + lbl_w + ci * cell_w
        add_rect(slide, cx, tbl_y, cell_w, hdr_h, fill_hex=primary)
        add_text(slide, dom[:12], cx + 0.03, tbl_y + 0.02, cell_w - 0.04, hdr_h - 0.04,
                 font_size=max(6.0, min(8.0, cell_w * 7)), bold=True, color_hex='FFFFFF',
                 align='center', wrap=False)

    # Rows
    for ri, row_dom in enumerate(doms):
        ry = tbl_y + hdr_h + ri * cell_h
        # Row label
        add_rect(slide, tbl_x, ry, lbl_w, cell_h, fill_hex='F1F5F9', line_hex='E2E8F0', line_w=0.3)
        add_text(slide, row_dom, tbl_x + 0.06, ry + 0.02, lbl_w - 0.08, cell_h - 0.04,
                 font_size=max(7.0, min(9.0, lbl_w * 5)), bold=True, color_hex='334155')
        # Cells
        for ci, col_dom in enumerate(doms):
            cx = tbl_x + lbl_w + ci * cell_w
            v = matrix[row_dom].get(col_dom, 0)
            if row_dom == col_dom:
                add_rect(slide, cx, ry, cell_w, cell_h, fill_hex='E2E8F0')
            else:
                intensity = int(v / max_v * 85) if v > 0 else 0
                bg = f"{255 - intensity:02X}{255 - intensity // 3:02X}FF" if v > 0 else 'FAFAFA'
                add_rect(slide, cx, ry, cell_w, cell_h, fill_hex=bg, line_hex='E2E8F0', line_w=0.3)
                if v > 0:
                    add_text(slide, str(v), cx, ry + 0.02, cell_w, cell_h - 0.04,
                             font_size=max(7.0, min(10.0, cell_w * 8)), bold=True,
                             color_hex=primary if v > 0 else '94A3B8', align='center')

    # Legend
    leg_y = tbl_y + hdr_h + n * cell_h + 0.15
    add_text(slide, "Les chiffres indiquent le nombre de flux inter-domaines",
             tbl_x, leg_y, CONTENT_W, 0.20, font_size=8, italic=True, color_hex='94A3B8')


def add_hubs_slide(prs: Presentation, apps: List[Dict], flows: List[Dict], opts: Dict):
    """Slide : top applications hub (les plus connectées)."""
    if not flows:
        return
    primary = opts.get('clientPrimary', '2979FF')
    layout = prs.slide_layouts[8]
    slide = prs.slides.add_slide(layout)

    # Connectivity per app
    conn = {a['id']: 0 for a in apps}
    flow_details = {a['id']: [] for a in apps}
    app_by_id = {a['id']: a for a in apps}
    for f in flows:
        fid, tid = f.get('from', ''), f.get('to', '')
        if fid in conn:
            conn[fid] += 1
            ta = app_by_id.get(tid, {})
            flow_details[fid].append(ta.get('name', ''))
        if tid in conn:
            conn[tid] += 1
            fa = app_by_id.get(fid, {})
            flow_details[tid].append(fa.get('name', ''))

    top_hubs = sorted([a for a in apps if conn.get(a['id'], 0) > 0],
                      key=lambda a: -conn.get(a['id'], 0))[:12]

    n_hubs = len(top_hubs)
    set_placeholders(slide,
        title="HUBS APPLICATIFS — APPLICATIONS LES PLUS CONNECTÉES",
        breadcrumb="Cartographie Applicative",
        subtitle=f"{n_hubs} applications · classement par nombre de connexions"
    )

    if not top_hubs:
        add_text(slide, "Aucun flux défini", CONTENT_X, CONTENT_Y, CONTENT_W, 0.40,
                 font_size=12, color_hex='94A3B8', align='center')
        return

    max_conn = conn.get(top_hubs[0]['id'], 1) or 1
    row_h = min(0.36, CONTENT_H / n_hubs)
    bar_max_w = CONTENT_W * 0.40
    name_w = CONTENT_W * 0.22
    dom_w = CONTENT_W * 0.16
    bar_x = CONTENT_X + name_w + dom_w + 0.10
    cnt_x = bar_x + bar_max_w + 0.08
    cnt_w = 0.60
    detail_x = cnt_x + cnt_w + 0.08
    detail_w = CONTENT_W - (detail_x - CONTENT_X) - 0.10

    # Header
    hdr_h = 0.24
    for lbl, x, w in [('Application', CONTENT_X, name_w), ('Domaine', CONTENT_X + name_w, dom_w),
                       ('Connexions', bar_x, bar_max_w + cnt_w + 0.10),
                       ('Applications connectées', detail_x, detail_w)]:
        add_rect(slide, x, CONTENT_Y, w, hdr_h, fill_hex=primary)
        add_text(slide, lbl, x + 0.04, CONTENT_Y + 0.03, w - 0.06, hdr_h - 0.04,
                 font_size=8, bold=True, color_hex='FFFFFF')

    for ri, app in enumerate(top_hubs):
        ry = CONTENT_Y + hdr_h + ri * row_h
        bg = 'F8FAFC' if ri % 2 == 0 else 'FFFFFF'
        app_conn = conn.get(app['id'], 0)
        bw = (app_conn / max_conn) * bar_max_w

        # App name
        add_rect(slide, CONTENT_X, ry, name_w, row_h, fill_hex=bg, line_hex='E2E8F0', line_w=0.3)
        add_text(slide, app.get('name', ''), CONTENT_X + 0.06, ry + 0.02,
                 name_w - 0.08, row_h - 0.04, font_size=8, bold=True, color_hex='1E293B')

        # Domain
        add_rect(slide, CONTENT_X + name_w, ry, dom_w, row_h, fill_hex=bg, line_hex='E2E8F0', line_w=0.3)
        add_text(slide, app.get('domain', ''), CONTENT_X + name_w + 0.04, ry + 0.02,
                 dom_w - 0.06, row_h - 0.04, font_size=7.5, color_hex='64748B')

        # Bar
        add_rect(slide, bar_x, ry, bar_max_w + cnt_w + 0.10, row_h, fill_hex=bg, line_hex='E2E8F0', line_w=0.3)
        if bw > 0:
            add_rect(slide, bar_x + 0.04, ry + (row_h - 0.14) / 2, bw, 0.14, fill_hex=primary)
        add_text(slide, str(app_conn), bar_x + 0.08 + bw, ry + 0.02,
                 cnt_w, row_h - 0.04, font_size=9, bold=True, color_hex=primary)

        # Connected apps
        add_rect(slide, detail_x, ry, detail_w, row_h, fill_hex=bg, line_hex='E2E8F0', line_w=0.3)
        details = ', '.join(flow_details[app['id']][:6])
        if len(flow_details[app['id']]) > 6:
            details += f" +{len(flow_details[app['id']]) - 6}"
        add_text(slide, details, detail_x + 0.04, ry + 0.02,
                 detail_w - 0.06, row_h - 0.04, font_size=7, color_hex='475569')


def add_env_applicatif_slides(prs: Presentation, apps: List[Dict], opts: Dict, field: str,
                               title: str, color_map: Dict):
    """Slides environnement applicatif — apps groupées par domaine avec statut coloré."""
    primary = opts.get('clientPrimary', '2979FF')

    doms = list(dict.fromkeys(a.get('domain', 'Autre') for a in apps))

    chip_w, chip_h = 1.30, 0.28
    chip_gap = 0.05
    dom_label_w = 1.60
    chips_area_w = CONTENT_W - dom_label_w - 0.10
    chips_per_row = max(1, int((chips_area_w + chip_gap) / (chip_w + chip_gap)))
    dom_pad = 0.08
    dom_gap = 0.08

    def make_slide():
        layout = prs.slide_layouts[8]
        sl = prs.slides.add_slide(layout)
        set_placeholders(sl, title=title, breadcrumb="Cartographie Applicative")
        # Legend
        lx = CONTENT_X
        ly = CONTENT_Y - 0.04
        for st, col in list(color_map.items())[:5]:
            add_rect(sl, lx, ly + 0.04, 0.14, 0.12, fill_hex=col)
            add_text(sl, st, lx + 0.18, ly, 1.30, 0.22, font_size=8, color_hex='334155')
            lx += 1.55
        return sl

    cur_slide = make_slide()
    cur_y = CONTENT_Y + 0.26
    max_y = SLIDE_H - 0.40

    for dom in doms:
        dom_apps = [a for a in apps if a.get('domain') == dom]
        chip_rows = math.ceil(len(dom_apps) / chips_per_row) or 1
        dom_h = dom_pad * 2 + chip_rows * chip_h + (chip_rows - 1) * chip_gap

        if cur_y + dom_h > max_y:
            cur_slide = make_slide()
            cur_y = CONTENT_Y + 0.26

        # Domain label box
        add_rect(cur_slide, CONTENT_X, cur_y, dom_label_w, dom_h,
                 fill_hex='EFF6FF', line_hex='BFDBFE', line_w=0.5)
        add_text(cur_slide, dom, CONTENT_X + 0.08, cur_y + dom_pad,
                 dom_label_w - 0.12, dom_h * 0.55, font_size=9, bold=True, color_hex='1D4ED8')
        add_text(cur_slide, f"{len(dom_apps)} app{'s' if len(dom_apps)>1 else ''}",
                 CONTENT_X + 0.08, cur_y + dom_h * 0.62,
                 dom_label_w - 0.12, dom_h * 0.32, font_size=7, color_hex='64748B')

        chip_x0 = CONTENT_X + dom_label_w + 0.10
        for ai, app in enumerate(dom_apps):
            row = ai // chips_per_row
            col = ai % chips_per_row
            cx = chip_x0 + col * (chip_w + chip_gap)
            cy = cur_y + dom_pad + row * (chip_h + chip_gap)
            st = app.get(field) or 'Non défini'
            stc = color_map.get(st, '94A3B8')
            add_rect(cur_slide, cx, cy, chip_w, chip_h,
                     fill_hex=stc, transparency=80, line_hex=stc, line_w=0.5)
            add_rect(cur_slide, cx, cy, 0.05, chip_h, fill_hex=stc)
            add_text(cur_slide, app.get('name', ''), cx + 0.09, cy + 0.01,
                     chip_w - 0.12, chip_h * 0.55, font_size=7, bold=True, color_hex='1a1a1a')
            add_text(cur_slide, '—' if st == 'Non défini' else st,
                     cx + 0.09, cy + chip_h * 0.55, chip_w - 0.12, chip_h * 0.42,
                     font_size=6.5, color_hex=stc)

        add_rect(cur_slide, CONTENT_X, cur_y + dom_h, CONTENT_W, 0.01, fill_hex='E2E8F0')
        cur_y += dom_h + dom_gap


def add_fin_slide(prs: Presentation):
    """Final slide."""
    layout = prs.slide_layouts[14]  # 6_Slide de fin
    slide = prs.slides.add_slide(layout)
    today = datetime.date.today().strftime("%d/%m/%Y")
    for ph in slide.placeholders:
        name = ph.name.lower()
        try:
            if 'date' in name:
                ph.text = today
        except Exception:
            pass


# ─── MAIN ENTRY POINT ────────────────────────────────────────────────────────
def generate_pptx(apps: List[Dict], flows: List[Dict], opts: Dict, template_path: str) -> bytes:
    """
    Generate a PPTX presentation using the given template.
    Returns the file as bytes.
    """
    prs = Presentation(template_path)
    _clear_slides(prs)

    incl_exec = opts.get('inclExecSlides', True)
    incl_paysage = opts.get('inclPaysage', True)
    incl_matrices = opts.get('inclMatrices', True)
    incl_domain_status = opts.get('inclDomainStatus', False)
    incl_carto = opts.get('inclConsolidatedCarto', True)

    if incl_exec:
        add_cover_slide(prs, apps, flows, opts)
        add_synthese_messages_slide(prs, apps, flows, opts)
        add_executive_slide(prs, apps, opts)

    if incl_carto and flows:
        add_flux_matrix_slide(prs, apps, flows, opts)
        add_hubs_slide(prs, apps, flows, opts)

    if incl_domain_status:
        add_domain_status_slides(prs, apps, opts)

    # Environnement applicatif D1 & D2 (toujours inclus si apps)
    if apps:
        add_env_applicatif_slides(prs, apps, opts, 'statusD1',
            "ENVIRONNEMENT APPLICATIF — DAY 1 (CLOSING)", D1_COLORS)
        add_env_applicatif_slides(prs, apps, opts, 'statusD2',
            "ENVIRONNEMENT APPLICATIF — DAY 2 (CIBLE)", D2_COLORS)

    if incl_paysage and apps:
        add_paysage_slide(prs, apps)

    if incl_matrices:
        add_recap_apps_slides(prs, apps, opts)

    add_flux_table_slides(prs, apps, flows, opts)
    add_fin_slide(prs)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()
